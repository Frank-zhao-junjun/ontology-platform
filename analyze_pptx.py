import sys
import json
sys.path.insert(0, 'C:/Users/Frank/.workbuddy/binaries/python/versions/3.13.12/Lib/site-packages')
from pptx import Presentation
from pptx.util import Pt, Emu

pptx_path = 'E:/00 - AI/本体建模/本体建模平台-产品介绍.pptx'
prs = Presentation(pptx_path)

result = {
    'slide_width_inches': prs.slide_width / 914400,
    'slide_height_inches': prs.slide_height / 914400,
    'total_slides': len(prs.slides),
    'slides': []
}

all_colors = set()
shape_types = {}

for idx, slide in enumerate(prs.slides, 1):
    slide_data = {
        'slide_number': idx,
        'layout_name': slide.slide_layout.name,
        'shapes': []
    }
    
    # Background
    try:
        bg = slide.background
        if bg.fill.type is not None:
            slide_data['background_fill_type'] = str(bg.fill.type)
    except:
        pass
    
    for shape in slide.shapes:
        shape_info = {
            'name': shape.name,
            'shape_type': str(shape.shape_type),
            'left': shape.left,
            'top': shape.top,
            'width': shape.width,
            'height': shape.height,
            'left_inches': round(shape.left / 914400, 2),
            'top_inches': round(shape.top / 914400, 2),
            'width_inches': round(shape.width / 914400, 2),
            'height_inches': round(shape.height / 914400, 2),
        }
        
        # Track shape types
        st = str(shape.shape_type)
        shape_types[st] = shape_types.get(st, 0) + 1
        
        # Fill color
        try:
            if hasattr(shape, 'fill') and shape.fill.type is not None:
                try:
                    color = shape.fill.fore_color.rgb
                    shape_info['fill_color'] = '#' + str(color)
                    all_colors.add(str(color))
                except:
                    pass
        except:
            pass
        
        # Line/border
        try:
            if hasattr(shape, 'line'):
                line = shape.line
                if line.fill.type is not None:
                    try:
                        line_color = line.fill.fore_color.rgb
                        shape_info['line_color'] = '#' + str(line_color)
                        all_colors.add(str(line_color))
                    except:
                        pass
        except:
            pass
        
        # Text frames
        if shape.has_text_frame:
            tf = shape.text_frame
            text_data = {
                'word_wrap': tf.word_wrap,
                'paragraphs': []
            }
            
            for para in tf.paragraphs:
                para_text = para.text
                if para_text.strip():
                    para_data = {
                        'text': para_text,
                        'alignment': str(para.alignment) if para.alignment is not None else None,
                        'level': para.level,
                        'runs': []
                    }
                    
                    for run in para.runs:
                        if run.text.strip():
                            run_data = {'text': run.text}
                            font = run.font
                            try:
                                if font.name:
                                    run_data['font'] = font.name
                                if font.size:
                                    run_data['size_pt'] = round(font.size / 12700, 1)
                                if font.bold:
                                    run_data['bold'] = True
                                if font.italic:
                                    run_data['italic'] = True
                                try:
                                    if font.color.rgb:
                                        c = str(font.color.rgb)
                                        run_data['color'] = '#' + c
                                        all_colors.add(c)
                                except:
                                    pass
                            except:
                                pass
                            para_data['runs'].append(run_data)
                    
                    text_data['paragraphs'].append(para_data)
            
            shape_info['text_frame'] = text_data
        
        # Images
        if shape.shape_type == 13:
            shape_info['is_image'] = True
            try:
                img = shape.image
                shape_info['image_type'] = img.content_type
            except:
                pass
        
        # Tables
        if shape.has_table:
            table = shape.table
            table_data = {
                'rows': len(table.rows),
                'cols': len(table.columns),
                'cells': []
            }
            for row_idx, row in enumerate(table.rows):
                for col_idx, cell in enumerate(row.cells):
                    if cell.text.strip():
                        table_data['cells'].append({
                            'row': row_idx,
                            'col': col_idx,
                            'text': cell.text
                        })
            shape_info['table'] = table_data
        
        slide_data['shapes'].append(shape_info)
    
    result['slides'].append(slide_data)

result['all_colors'] = sorted(['#' + c for c in all_colors])

# Write to JSON file (UTF-8)
output_path = 'E:/00 - AI/本体建模/pptx_analysis.json'
with open(output_path, 'w', encoding='utf-8') as f:
    json.dump(result, f, ensure_ascii=False, indent=2)

print(f"Analysis written to: {output_path}")
print(f"Total slides: {len(result['slides'])}")
print(f"Total colors: {len(all_colors)}")
print(f"Shape types: {shape_types}")
print(f"Colors: {sorted(['#' + c for c in all_colors])}")
